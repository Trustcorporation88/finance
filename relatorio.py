# -*- coding: utf-8 -*-
"""Geração de relatórios em Word (.docx), PDF e PowerPoint (.pptx).

- Word: python-docx, com os gráficos embutidos.
- PDF: reportlab (platypus), com os gráficos embutidos.
- PPTX: python-pptx, com os gráficos embutidos.
"""
from __future__ import annotations

import io
import base64

import os

import markdown_utils as mdu

# Marca TRUST
LOGO_TRUST = os.path.join(os.path.dirname(__file__), "static", "logo-trust.png")
MARCA = "TRUST CORPORATION · CFO de Bolso"

from docx import Document
from docx.shared import Pt, Inches, RGBColor

try:
    from pptx import Presentation as PptxPresentation
    from pptx.util import Inches as PptxInches, Pt as PptxPt
    from pptx.dml.color import RGBColor as PptxRGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    PPTX_DISPONIVEL = True
except ImportError:
    PPTX_DISPONIVEL = False

from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import mm
from reportlab.lib.colors import HexColor
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle
from reportlab.lib.enums import TA_CENTER

AZUL = RGBColor(0x1F, 0x3B, 0x73)
VERDE = RGBColor(0x1E, 0x8E, 0x3E)
VERMELHO = RGBColor(0xC0, 0x39, 0x2B)
CINZA = RGBColor(0x55, 0x55, 0x55)
AMARELO = RGBColor(0xF3, 0x9C, 0x12)


def _b64_para_bytes(b64: str) -> bytes:
    return base64.b64decode(b64.split(",", 1)[1] if "," in b64 else b64)


# ==========================================================================
# WORD
# ==========================================================================

def gerar_word(resultado: dict, graf: dict) -> io.BytesIO:
    doc = Document()
    estilo = doc.styles["Normal"]
    estilo.font.name = "Calibri"
    estilo.font.size = Pt(11)

    doc.add_heading("Raio-X do Caixa — CFO de Bolso", level=0)

    p = doc.add_paragraph()
    r = p.add_run("Números calculados em regime de caixa, direto dos seus dados.")
    r.italic = True
    r.font.size = Pt(9)
    r.font.color.rgb = CINZA

    doc.add_heading("O essencial (em 1 olhada)", level=1)
    t = doc.add_table(rows=4, cols=2)
    t.style = "Light Grid Accent 1"
    linhas = [
        ("Entrou", f"R$ {resultado.get('total_entradas', 0):,.2f}"),
        ("Saiu", f"R$ {resultado.get('total_saidas', 0):,.2f}"),
        ("Sobrou / Faltou", f"R$ {resultado.get('saldo', 0):,.2f}"),
        ("Margem do período", f"{resultado.get('margem')}%" if resultado.get("margem") is not None else "não calculável"),
    ]
    for i, (k, v) in enumerate(linhas):
        t.cell(i, 0).text = k
        t.cell(i, 1).text = v
    _estilizar_tabela_docx(t)

    if graf.get("comparativo"):
        doc.add_picture(io.BytesIO(_b64_para_bytes(graf["comparativo"])), width=Inches(5.5))
        doc.add_paragraph()

    if resultado.get("gastos_categorias"):
        doc.add_heading("Pra onde foi o dinheiro", level=1)
        total_s = resultado.get("total_saidas", 1) or 1
        linhas_cat = [("Categoria", "Valor", "%")]
        for cat, val in resultado["gastos_categorias"].items():
            linhas_cat.append((cat, f"R$ {val:,.2f}", f"{val / total_s * 100:.1f}%"))
        t = doc.add_table(rows=len(linhas_cat), cols=3)
        t.style = "Light Grid Accent 1"
        for i, linha in enumerate(linhas_cat):
            for j, v in enumerate(linha):
                t.cell(i, j).text = v
        _estilizar_tabela_docx(t)
        if graf.get("pizza_gastos"):
            doc.add_picture(io.BytesIO(_b64_para_bytes(graf["pizza_gastos"])), width=Inches(5))
            doc.add_paragraph()

    if graf.get("fluxo_mensal"):
        doc.add_heading("Fluxo de caixa por mês", level=1)
        doc.add_picture(io.BytesIO(_b64_para_bytes(graf["fluxo_mensal"])), width=Inches(5.5))
        doc.add_paragraph()

    if resultado.get("alertas"):
        doc.add_heading("Alertas", level=1)
        for a in resultado["alertas"]:
            cor = VERMELHO if a["nivel"] == "alto" else AMARELO
            p = doc.add_paragraph()
            r = p.add_run(f"[{a['nivel'].upper()}] {a['titulo']} — ")
            r.bold = True
            r.font.color.rgb = cor
            p.add_run(a["detalhe"])

    if resultado.get("narrativa_ia"):
        doc.add_heading("Análise da IA (DeepSeek)", level=1)
        _adicionar_blocos_word(doc, mdu.parse_blocks(resultado["narrativa_ia"]))
    elif resultado.get("erro_ia"):
        doc.add_heading("IA indisponível", level=1)
        p = doc.add_paragraph()
        r = p.add_run("Não foi possível consultar a IA: ")
        r.bold = True
        p.add_run(resultado["erro_ia"])

    doc.add_paragraph()
    p = doc.add_paragraph()
    r = p.add_run("Lembrete: é gestão de caixa gerencial. Para imposto/nota fiscal, fale com seu contador.")
    r.italic = True
    r.font.size = Pt(8)
    r.font.color.rgb = CINZA

    buf = io.BytesIO()
    # rodapé da marca
    _rodape_word(doc)
    doc.save(buf)
    buf.seek(0)
    return buf


def _estilizar_tabela_docx(t):
    for i, row in enumerate(t.rows):
        for cell in row.cells:
            for pr in cell.paragraphs:
                for r in pr.runs:
                    r.font.size = Pt(10)
                    if i == 0:
                        r.bold = True


def _adicionar_runs_docx(p, texto_marcado):
    """Adiciona runs com negrito/itálico num parágrafo docx."""
    for txt, bold, italic in mdu.runs_do_texto(texto_marcado):
        r = p.add_run(txt)
        r.bold = bold
        r.italic = italic
        r.font.size = Pt(11)


def _adicionar_blocos_word(doc, blocos):
    """Renderiza blocos de markdown num documento Word."""
    for b in blocos:
        tipo = b["tipo"]
        if tipo == "hr":
            doc.add_paragraph("─" * 50)
        elif tipo == "p":
            p = doc.add_paragraph()
            _adicionar_runs_docx(p, b["texto"])
        elif tipo in ("h1", "h2", "h3", "h4"):
            nivel = int(tipo[1]) - 1
            doc.add_heading(b["texto"], level=nivel)
        elif tipo in ("ul", "ol"):
            for item in b["itens"]:
                p = doc.add_paragraph(style="List Bullet" if tipo == "ul" else "List Number")
                _adicionar_runs_docx(p, item)
        elif tipo == "blockquote":
            p = doc.add_paragraph()
            _adicionar_runs_docx(p, b["texto"])
            for pr in p.runs:
                pr.italic = True
                pr.font.color.rgb = CINZA
        elif tipo == "table":
            linhas = b["linhas"]
            t = doc.add_table(rows=len(linhas), cols=len(linhas[0]))
            t.style = "Light Grid Accent 1"
            for i, linha in enumerate(linhas):
                for j, cel in enumerate(linha):
                    if j < len(linhas[0]):
                        t.cell(i, j).text = cel
            _estilizar_tabela_docx(t)
            doc.add_paragraph()


AMARELO_RL = HexColor("#F39C12")


# ==========================================================================
# PDF
# ==========================================================================

def gerar_pdf(resultado: dict, graf: dict) -> io.BytesIO:
    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4,
                            leftMargin=15 * mm, rightMargin=15 * mm,
                            topMargin=15 * mm, bottomMargin=15 * mm)

    estilos = getSampleStyleSheet()
    titulo = ParagraphStyle("Titulo", parent=estilos["Title"], fontSize=20,
                            textColor=HexColor("#1F3B73"), spaceAfter=4)
    h2 = ParagraphStyle("H2", parent=estilos["Heading2"], fontSize=14,
                        textColor=HexColor("#1F3B73"), spaceBefore=12, spaceAfter=4)
    corpo = ParagraphStyle("Corpo", parent=estilos["BodyText"], fontSize=10, leading=14)
    nota = ParagraphStyle("Nota", parent=estilos["BodyText"], fontSize=8,
                          textColor=HexColor("#555555"), italic=True)

    elementos = []
    _marca_pdf(doc, elementos)
    elementos.append(Paragraph("Raio-X do Caixa \u2014 CFO de Bolso", titulo))
    elementos.append(Paragraph("N\u00fameros calculados em regime de caixa, direto dos seus dados.", nota))
    elementos.append(Spacer(1, 8))

    elementos.append(Paragraph("O essencial (em 1 olhada)", h2))
    linhas = [
        ["Entrou", f"R$ {resultado.get('total_entradas', 0):,.2f}"],
        ["Saiu", f"R$ {resultado.get('total_saidas', 0):,.2f}"],
        ["Sobrou / Faltou", f"R$ {resultado.get('saldo', 0):,.2f}"],
        ["Margem do período", f"{resultado.get('margem')}%" if resultado.get("margem") is not None else "não calculável"],
    ]
    tbl = Table(linhas, colWidths=[80 * mm, 90 * mm])
    tbl.setStyle(TableStyle([
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTSIZE", (0, 0), (-1, -1), 10),
        ("GRID", (0, 0), (-1, -1), 0.5, HexColor("#BBBBBB")),
        ("BACKGROUND", (0, 0), (0, -1), HexColor("#E8EDF4")),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
    ]))
    elementos.append(tbl)
    elementos.append(Spacer(1, 8))

    if graf.get("comparativo"):
        img = Image(io.BytesIO(_b64_para_bytes(graf["comparativo"])), width=150 * mm, height=90 * mm)
        img.hAlign = "CENTER"
        elementos.append(img)

    if resultado.get("gastos_categorias"):
        elementos.append(Paragraph("Pra onde foi o dinheiro", h2))
        total_s = resultado.get("total_saidas", 1) or 1
        linhas_cat = [["Categoria", "Valor", "%"]]
        for cat, val in resultado["gastos_categorias"].items():
            linhas_cat.append([cat, f"R$ {val:,.2f}", f"{val / total_s * 100:.1f}%"])
        tbl = Table(linhas_cat, colWidths=[90 * mm, 45 * mm, 35 * mm])
        tbl.setStyle(TableStyle([
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, -1), 9),
            ("GRID", (0, 0), (-1, -1), 0.4, HexColor("#CCCCCC")),
            ("BACKGROUND", (0, 0), (-1, 0), HexColor("#1F3B73")),
            ("TEXTCOLOR", (0, 0), (-1, 0), HexColor("#FFFFFF")),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ]))
        elementos.append(tbl)
        elementos.append(Spacer(1, 8))
        if graf.get("pizza_gastos"):
            img = Image(io.BytesIO(_b64_para_bytes(graf["pizza_gastos"])), width=120 * mm, height=90 * mm)
            img.hAlign = "CENTER"
            elementos.append(img)

    if graf.get("fluxo_mensal"):
        elementos.append(Paragraph("Fluxo de caixa por mês", h2))
        img = Image(io.BytesIO(_b64_para_bytes(graf["fluxo_mensal"])), width=150 * mm, height=90 * mm)
        img.hAlign = "CENTER"
        elementos.append(img)

    if resultado.get("alertas"):
        elementos.append(Paragraph("Alertas", h2))
        for a in resultado["alertas"]:
            cor = HexColor("#C0392B") if a["nivel"] == "alto" else AMARELO_RL
            elementos.append(Paragraph(f"<font color='{cor.hexval()}'><b>[{a['nivel'].upper()}] {a['titulo']}</b></font> — {a['detalhe']}", corpo))

    if resultado.get("narrativa_ia"):
        elementos.append(Paragraph("Análise da IA (DeepSeek)", h2))
        _adicionar_blocos_pdf(elementos, mdu.parse_blocks(resultado["narrativa_ia"]), corpo, h2)

    elementos.append(Spacer(1, 12))
    elementos.append(Paragraph("Lembrete: é gestão de caixa gerencial. Para imposto/nota fiscal, fale com seu contador.", nota))
    _rodape_pdf(elementos)

    doc.build(elementos)
    buf.seek(0)
    return buf


def _adicionar_blocos_pdf(elementos, blocos, corpo, h2):
    """Renderiza blocos de markdown em elementos do PDF (reportlab platypus)."""
    for b in blocos:
        tipo = b["tipo"]
        if tipo == "hr":
            from reportlab.platypus import HRFlowable
            elementos.append(HRFlowable(width="100%", thickness=0.7, color=HexColor("#BBBBBB")))
            elementos.append(Spacer(1, 4))
        elif tipo == "p":
            texto = mdu.runs_para_reportlab(b["texto"])
            elementos.append(Paragraph(texto, corpo))
            elementos.append(Spacer(1, 4))
        elif tipo in ("h1", "h2", "h3", "h4"):
            st = ParagraphStyle(
                f"HIA{tipo}", parent=h2,
                fontSize=max(9, 15 - 1.6 * int(tipo[1])),
                spaceBefore=8, spaceAfter=3,
            )
            elementos.append(Paragraph(mdu.runs_para_reportlab(b["texto"]), st))
        elif tipo in ("ul", "ol"):
            for item in b["itens"]:
                marc = "• " if tipo == "ul" else "• "
                elementos.append(Paragraph(marc + mdu.runs_para_reportlab(item), corpo))
                elementos.append(Spacer(1, 2))
        elif tipo == "blockquote":
            st = ParagraphStyle("BQIA", parent=corpo, fontStyle="italic",
                                textColor=HexColor("#555555"), leftIndent=12)
            elementos.append(Paragraph(mdu.runs_para_reportlab(b["texto"]), st))
            elementos.append(Spacer(1, 4))
        elif tipo == "table":
            linhas = b["linhas"]
            ncol = max(len(r) for r in linhas)
            linhas = [r + [""] * (ncol - len(r)) for r in linhas]
            larg = 175 / ncol if ncol else 175
            tbl = Table(linhas, colWidths=[larg * mm] * ncol)
            tbl.setStyle(TableStyle([
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTSIZE", (0, 0), (-1, -1), 8.5),
                ("GRID", (0, 0), (-1, -1), 0.4, HexColor("#CCCCCC")),
                ("BACKGROUND", (0, 0), (-1, 0), HexColor("#1F3B73")),
                ("TEXTCOLOR", (0, 0), (-1, 0), HexColor("#FFFFFF")),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [HexColor("#FFFFFF"), HexColor("#F2F5F9")]),
                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
            ]))
            elementos.append(tbl)
            elementos.append(Spacer(1, 6))
    """Adiciona o rodapé com a marca TRUST no fim do documento."""
    try:
        doc.add_paragraph()
        p = doc.add_paragraph()
        r = p.add_run(MARCA)
        r.italic = True
        r.font.size = Pt(8)
        r.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    except Exception:
        pass


def _marca_pdf(doc, elementos):
    """Adiciona logotipo e rodapé da marca TRUST no PDF."""
    from reportlab.platypus import Image as RLImage, HRFlowable
    try:
        if os.path.exists(LOGO_TRUST):
            img = RLImage(LOGO_TRUST, width=45 * mm, height=45 * mm)
            img.hAlign = "LEFT"
            elementos.append(img)
        else:
            elementos.append(Paragraph("<b>TRUST CORPORATION</b>", estilos_titulo()))
        elementos.append(HRFlowable(width="100%", thickness=1, color=HexColor("#1F3B73")))
    except Exception:
        pass


def estilos_titulo():
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.colors import HexColor
    e = getSampleStyleSheet()
    return ParagraphStyle("TituloTrust", parent=e["Title"], fontSize=18,
                          textColor=HexColor("#1F3B73"))


def _rodape_word(doc):
    """Adiciona o rodapé com a marca TRUST no fim do documento."""
    try:
        doc.add_paragraph()
        p = doc.add_paragraph()
        r = p.add_run(MARCA)
        r.italic = True
        r.font.size = Pt(8)
        r.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    except Exception:
        pass


def _rodape_pdf(elementos, texto="TRUST CORPORATION"):
    from reportlab.platypus import Spacer, Paragraph
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.colors import HexColor
    st = ParagraphStyle("Rodape", fontSize=8, textColor=HexColor("#888888"), alignment=1)
    elementos.append(Spacer(1, 10))
    elementos.append(Paragraph(texto, st))


# ==========================================================================
# POWERPOINT
# ==========================================================================

AZUL_P = PptxRGBColor(0x3E, 0x7C, 0x3E)
VERDE_P = PptxRGBColor(0x1E, 0x8E, 0x3E)
VERMELHO_P = PptxRGBColor(0xC0, 0x39, 0x2B)
BRANCO_P = PptxRGBColor(0xFF, 0xFF, 0xFF)
CINZA_P = PptxRGBColor(0x55, 0x55, 0x55)
AMARELO_P = PptxRGBColor(0xF3, 0x9C, 0x12)
CLARO_P = PptxRGBColor(0xE8, 0xF2, 0xE8)


def gerar_pptx(resultado: dict, graf: dict) -> io.BytesIO:
    """Gera uma apresentação .pptx com o raio-x do caixa."""
    if not PPTX_DISPONIVEL:
        raise RuntimeError("python-pptx não instalado no servidor.")

    prs = PptxPresentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)
    blank = prs.slide_layouts[6]

    def novo_slide():
        return prs.slides.add_slide(blank)

    def caixa(slide, l, t, w, h, cor=None):
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptxInches(l), PptxInches(t), PptxInches(w), PptxInches(h))
        if cor is None:
            sh.fill.background()
        else:
            sh.fill.solid()
            sh.fill.fore_color.rgb = cor
        sh.line.fill.background()
        return sh

    def texto(slide, l, t, w, h, conteudo, size=18, bold=False, cor=CINZA_P, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(PptxInches(l), PptxInches(t), PptxInches(w), PptxInches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = conteudo
        r.font.size = PptxPt(size)
        r.font.bold = bold
        r.font.color.rgb = cor
        return tb

    def cabecalho(slide, titulo, sub=None):
        caixa(slide, 0, 0, 13.333, 0.95, AZUL_P)
        try:
            if os.path.exists(LOGO_TRUST):
                slide.shapes.add_picture(LOGO_TRUST, PptxInches(0.3), PptxInches(0.14), height=PptxInches(0.65))
                texto(slide, 2.3, 0.16, 10.9, 0.7, titulo, size=24, bold=True, cor=BRANCO_P)
                if sub:
                    texto(slide, 2.3, 0.55, 10.9, 0.4, sub, size=12, cor=PptxRGBColor(0xE8, 0xF2, 0xE8))
            else:
                texto(slide, 0.4, 0.16, 12.5, 0.7, titulo, size=24, bold=True, cor=BRANCO_P)
                if sub:
                    texto(slide, 0.4, 0.55, 12.5, 0.4, sub, size=12, cor=PptxRGBColor(0xE8, 0xF2, 0xE8))
        except Exception:
            texto(slide, 0.4, 0.16, 12.5, 0.7, titulo, size=24, bold=True, cor=BRANCO_P)
            if sub:
                texto(slide, 0.4, 0.55, 12.5, 0.4, sub, size=12, cor=PptxRGBColor(0xE8, 0xF2, 0xE8))

    def tabela(slide, l, t, w, rows, col_w=None, font_size=12, header_font=12):
        nrows = len(rows)
        ncols = len(rows[0])
        shp = slide.shapes.add_table(nrows, ncols, PptxInches(l), PptxInches(t), PptxInches(w), PptxInches(0.32 * nrows))
        tbl = shp.table
        if col_w:
            for j, cw in enumerate(col_w):
                tbl.columns[j].width = PptxInches(cw)
        for i, row in enumerate(rows):
            for j, val in enumerate(row):
                cell = tbl.cell(i, j)
                cell.text = str(val)
                for p in cell.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.size = PptxPt(header_font if i == 0 else font_size)
                        r.font.bold = (i == 0)
                        r.font.color.rgb = BRANCO_P if i == 0 else PptxRGBColor(0x22, 0x22, 0x22)
                if i == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = AZUL_P
        return tbl

    def add_pic(slide, b64, l, t, w):
        slide.shapes.add_picture(io.BytesIO(_b64_para_bytes(b64)), PptxInches(l), PptxInches(t), width=PptxInches(w))

    def add_linhas(slide, l, t, w, h, linhas, size=14):
        tb = slide.shapes.add_textbox(PptxInches(l), PptxInches(t), PptxInches(w), PptxInches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, (txt_, bold, cor, sz) in enumerate(linhas):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            r = p.add_run()
            r.text = txt_
            r.font.size = PptxPt(sz)
            r.font.bold = bold
            r.font.color.rgb = cor
            p.space_after = PptxPt(4)

    # ---- Slide 1: capa ----
    s = novo_slide()
    caixa(s, 0, 0, 13.333, 7.5, AZUL_P)
    caixa(s, 0, 4.5, 13.333, 0.06, AMARELO_P)
    texto(s, 1.0, 2.4, 11.3, 1.0, "RAIO-X DO CAIXA", size=44, bold=True, cor=BRANCO_P)
    texto(s, 1.0, 3.4, 11.3, 0.7, "Análise automática dos seus dados", size=22, cor=PptxRGBColor(0xDD, 0xFF, 0xDD))

    # ---- Slide 2: essencial ----
    s = novo_slide()
    cabecalho(s, "O essencial (em 1 olhada)")
    margem = f"{resultado.get('margem')}%" if resultado.get("margem") is not None else "não calculável"
    tabela(s, 0.8, 1.6, 11.0, [
        ["Entrou", "Saiu", "Sobrou / Faltou", "Margem"],
        [f"R$ {resultado.get('total_entradas', 0):,.2f}", f"R$ {resultado.get('total_saidas', 0):,.2f}",
         f"R$ {resultado.get('saldo', 0):,.2f}", margem],
    ], col_w=[2.75, 2.75, 2.75, 2.75], font_size=13, header_font=12)
    if graf.get("comparativo"):
        add_pic(s, graf["comparativo"], 2.0, 3.4, 9.0)

    # ---- Slide 3: pra onde foi ----
    s = novo_slide()
    cabecalho(s, "Pra onde foi o dinheiro")
    total_s = resultado.get("total_saidas", 1) or 1
    cats = list(resultado.get("gastos_categorias", {}).items())
    if cats:
        metade = (len(cats) + 1) // 2
        esq = cats[:metade]
        dir_ = cats[metade:]
        linhas_esq = [["Categoria", "R$", "%"]] + [
            [c, f"{v:,.2f}", f"{v / total_s * 100:.1f}%"] for c, v in esq
        ]
        linhas_dir = [["Categoria", "R$", "%"]] + [
            [c, f"{v:,.2f}", f"{v / total_s * 100:.1f}%"] for c, v in dir_
        ]
        tabela(s, 0.4, 1.4, 6.0, linhas_esq, col_w=[3.2, 1.7, 1.1], font_size=10, header_font=10)
        tabela(s, 6.9, 1.4, 6.0, linhas_dir, col_w=[3.2, 1.7, 1.1], font_size=10, header_font=10)
    if graf.get("pizza_gastos"):
        add_pic(s, graf["pizza_gastos"], 3.3, 5.0, 6.5)

    # ---- Slide 4: fluxo mensal ----
    if graf.get("fluxo_mensal"):
        s = novo_slide()
        cabecalho(s, "Fluxo de caixa por mês")
        add_pic(s, graf["fluxo_mensal"], 1.7, 1.6, 10.0)
        if resultado.get("por_mes"):
            linhas = [["Mês", "Entradas", "Saídas", "Saldo"]]
            for mes, v in resultado["por_mes"].items():
                linhas.append([mes, f"R$ {v['entradas']:,.2f}", f"R$ {v['saidas']:,.2f}", f"R$ {v['saldo']:,.2f}"])
            tabela(s, 1.5, 5.6, 10.0, linhas, col_w=[2.5, 2.5, 2.5, 2.5], font_size=11, header_font=11)

    # ---- Slide 5: alertas ----
    s = novo_slide()
    cabecalho(s, "Alertas")
    alertas = resultado.get("alertas", [])
    linhas = []
    for a in alertas:
        cor = VERMELHO_P if a["nivel"] == "alto" else AMARELO_P
        linhas.append((f"[{a['nivel'].upper()}] {a['titulo']}", True, cor, 18))
        linhas.append((a["detalhe"], False, CINZA_P, 14))
        linhas.append(("", False, CINZA_P, 6))
    if not linhas:
        linhas = [("Nenhum alerta automático detectado.", False, VERDE_P, 16)]
    add_linhas(s, 0.7, 1.5, 12.0, 5.5, linhas, size=14)

    # ---- Slide 6: narrativa IA ----
    if resultado.get("narrativa_ia"):
        s = novo_slide()
        cabecalho(s, "Análise da IA (DeepSeek)")
        _adicionar_blocos_pptx(s, resultado["narrativa_ia"])

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ==========================================================================
# PLANILHA COMPLETA: geração de apresentação, relatório e exportação
# ==========================================================================

def _abas_para_linhas(info):
    """Converte abas da planilha em linhas de tabela para relatórios.
    Aceita abas como dicts (nome, cabecalhos, amostra) ou como strings (só nome)."""
    linhas = []
    for a in info.get("abas", []):
        if isinstance(a, dict):
            nome = a.get("nome", "")
            nlin = a.get("linhas_lidas", 0)
            ncol = a.get("colunas_max", 0)
            cab = "; ".join(" | ".join(str(x) for x in c) for c in a.get("cabecalhos", [])[:2])
        else:
            nome = str(a)
            nlin = ncol = 0
            cab = ""
        linhas.append([nome, f"{nlin} linhas", f"{ncol} col", cab[:80]])
    return linhas


def _adicionar_blocos_pptx(slide, md, size=14):
    """Renderiza blocos de markdown num slide do PPTX (texto simples formatado)."""
    def caixa_txt(l, t, w, h, conteudo, sz=size, bold=False, cor=CINZA_P):
        tb = slide.shapes.add_textbox(PptxInches(l), PptxInches(t), PptxInches(w), PptxInches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = conteudo
        r.font.size = PptxPt(sz)
        r.font.bold = bold
        r.font.color.rgb = cor
        return tb

    blocos = mdu.parse_blocks(md)
    y = 1.3
    for b in blocos:
        tipo = b["tipo"]
        if tipo in ("h1", "h2", "h3", "h4"):
            caixa_txt(0.6, y, 12.2, 0.5, b["texto"], sz=size + 4, bold=True, cor=AZUL_P)
            y += 0.5
        elif tipo == "p":
            caixa_txt(0.6, y, 12.2, 0.6, b["texto"], sz=size, cor=CINZA_P)
            y += 0.45
        elif tipo in ("ul", "ol"):
            for item in b["itens"]:
                caixa_txt(0.8, y, 11.8, 0.5, "•  " + item, sz=size, cor=CINZA_P)
                y += 0.4
        elif tipo == "blockquote":
            caixa_txt(0.8, y, 11.8, 0.5, "▸  " + b["texto"], sz=size, bold=True, cor=AMARELO_P)
            y += 0.4
        elif tipo == "table":
            linhas = b["linhas"]
            ncol = max(len(r) for r in linhas)
            linhas = [r + [""] * (ncol - len(r)) for r in linhas]
            # quebra a tabela em várias linhas de texto para não estourar o slide
            for linha in linhas:
                caixa_txt(0.6, y, 12.2, 0.4, "  ".join(str(c)[:22] for c in linha), sz=size - 2, cor=CINZA_P)
                y += 0.35
        y += 0.08
        if y > 6.9:
            break


def gerar_pptx_planilha(resultado: dict, info: dict = None) -> io.BytesIO:
    if not PPTX_DISPONIVEL:
        raise RuntimeError("python-pptx não instalado no servidor.")
    info = info or resultado.get("planilha_info", {})
    abas = info.get("abas", [])

    prs = PptxPresentation()
    prs.slide_width = PptxInches(13.333)
    prs.slide_height = PptxInches(7.5)
    blank = prs.slide_layouts[6]

    def novo_slide():
        return prs.slides.add_slide(blank)

    def caixa(slide, l, t, w, h, cor=None):
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptxInches(l), PptxInches(t), PptxInches(w), PptxInches(h))
        if cor is None:
            sh.fill.background()
        else:
            sh.fill.solid()
            sh.fill.fore_color.rgb = cor
        sh.line.fill.background()
        return sh

    def texto(slide, l, t, w, h, conteudo, size=18, bold=False, cor=CINZA_P, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(PptxInches(l), PptxInches(t), PptxInches(w), PptxInches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = conteudo
        r.font.size = PptxPt(size)
        r.font.bold = bold
        r.font.color.rgb = cor
        return tb

    def cabecalho(slide, titulo, sub=None):
        caixa(slide, 0, 0, 13.333, 0.95, AZUL_P)
        try:
            if os.path.exists(LOGO_TRUST):
                slide.shapes.add_picture(LOGO_TRUST, PptxInches(0.3), PptxInches(0.14), height=PptxInches(0.65))
                texto(slide, 2.3, 0.16, 10.9, 0.7, titulo, size=24, bold=True, cor=BRANCO_P)
                if sub:
                    texto(slide, 2.3, 0.55, 10.9, 0.4, sub, size=12, cor=PptxRGBColor(0xE8, 0xF2, 0xE8))
            else:
                texto(slide, 0.4, 0.16, 12.5, 0.7, titulo, size=24, bold=True, cor=BRANCO_P)
                if sub:
                    texto(slide, 0.4, 0.55, 12.5, 0.4, sub, size=12, cor=PptxRGBColor(0xE8, 0xF2, 0xE8))
        except Exception:
            texto(slide, 0.4, 0.16, 12.5, 0.7, titulo, size=24, bold=True, cor=BRANCO_P)
            if sub:
                texto(slide, 0.4, 0.55, 12.5, 0.4, sub, size=12, cor=PptxRGBColor(0xE8, 0xF2, 0xE8))

    def tabela(slide, l, t, w, rows, col_w=None, font_size=10, header_font=10):
        nrows = len(rows)
        ncols = len(rows[0])
        shp = slide.shapes.add_table(nrows, ncols, PptxInches(l), PptxInches(t), PptxInches(w), PptxInches(0.28 * nrows))
        tbl = shp.table
        if col_w:
            for j, cw in enumerate(col_w):
                tbl.columns[j].width = PptxInches(cw)
        for i, row in enumerate(rows):
            for j, val in enumerate(row):
                cell = tbl.cell(i, j)
                cell.text = str(val)
                for p in cell.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.size = PptxPt(header_font if i == 0 else font_size)
                        r.font.bold = (i == 0)
                        r.font.color.rgb = BRANCO_P if i == 0 else PptxRGBColor(0x22, 0x22, 0x22)
                if i == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = AZUL_P
        return tbl

    # Slide 1: capa
    s = novo_slide()
    caixa(s, 0, 0, 13.333, 7.5, AZUL_P)
    caixa(s, 0, 4.5, 13.333, 0.06, AMARELO_P)
    texto(s, 1.0, 2.3, 11.3, 1.0, "ANÁLISE DA PLANILHA", size=42, bold=True, cor=BRANCO_P)
    nome_arq = resultado.get("nome_arquivo", "planilha")
    texto(s, 1.0, 3.3, 11.3, 0.7, str(nome_arq), size=20, cor=PptxRGBColor(0xDD, 0xFF, 0xDD))
    texto(s, 1.0, 4.0, 11.3, 0.6, f"{info.get('n_abas', 0)} abas analisadas", size=16, cor=PptxRGBColor(0xCC, 0xFF, 0xCC))

    # Slide 2: visão geral das abas
    s = novo_slide()
    cabecalho(s, "Visão geral das abas")
    linhas_tab = [["Aba", "Linhas", "Colunas", "Conteúdo"]] + _abas_para_linhas(info)
    tabela(s, 0.4, 1.3, 12.5, linhas_tab, col_w=[3.2, 1.3, 1.2, 4.8], font_size=10, header_font=10)

    # Slides por aba (até 8 abas mais importantes para não estourar)
    slides_aba = abas[:8]
    for idx, a in enumerate(slides_aba):
        s = novo_slide()
        cabecalho(s, f"Aba: {a.get('nome','')}")
        cab = a.get("cabecalhos", [])
        amostra = a.get("amostra", [])
        y = 1.3
        if cab:
            texto(s, 0.4, y, 12.5, 0.4, "Conteúdo/cabeçalhos:", size=13, bold=True, cor=VERDE_P)
            y += 0.35
            for c in cab[:4]:
                texto(s, 0.6, y, 12.4, 0.35, "  " + " | ".join(str(x) for x in c)[:140], size=11, cor=CINZA_P)
                y += 0.32
        if amostra:
            texto(s, 0.4, y, 12.5, 0.4, "Dados:", size=13, bold=True, cor=VERDE_P)
            y += 0.35
            for c in amostra[:6]:
                texto(s, 0.6, y, 12.4, 0.32, "  " + " | ".join(str(x) for x in c)[:140], size=10, cor=CINZA_P)
                y += 0.30

    # Slide final: análise da IA
    if resultado.get("narrativa_ia"):
        s = novo_slide()
        cabecalho(s, "Análise da IA")
        _adicionar_blocos_pptx(s, resultado["narrativa_ia"])

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def gerar_word_planilha(resultado: dict, info: dict = None) -> io.BytesIO:
    """Gera um relatório Word resumindo a planilha completa."""
    info = info or resultado.get("planilha_info", {})
    doc = Document()
    estilo = doc.styles["Normal"]
    estilo.font.name = "Calibri"
    estilo.font.size = Pt(11)

    doc.add_heading("Análise da Planilha", level=0)
    nome = resultado.get("nome_arquivo", "planilha")
    p = doc.add_paragraph()
    r = p.add_run(f"Arquivo: {nome} · {info.get('n_abas', 0)} abas")
    r.italic = True
    r.font.size = Pt(9)
    r.font.color.rgb = CINZA

    # Tabela de abas
    doc.add_heading("Abas da planilha", level=1)
    linhas = [["Aba", "Linhas", "Colunas", "Conteúdo"]] + _abas_para_linhas(info)
    t = doc.add_table(rows=len(linhas), cols=4)
    t.style = "Light Grid Accent 1"
    for i, row in enumerate(linhas):
        for j, val in enumerate(row):
            t.cell(i, j).text = str(val)
    for i, row in enumerate(t.rows):
        for cell in row.cells:
            for pr in cell.paragraphs:
                for rr in pr.runs:
                    rr.font.size = Pt(9)
                    if i == 0:
                        rr.bold = True

    # Detalhe de cada aba
    for a in info.get("abas", []):
        doc.add_heading(f"Aba: {a.get('nome','')}", level=2)
        for c in a.get("cabecalhos", [])[:5]:
            doc.add_paragraph(" | ".join(str(x) for x in c))
        for c in a.get("amostra", [])[:8]:
            doc.add_paragraph("   " + " | ".join(str(x) for x in c))

    if resultado.get("narrativa_ia"):
        doc.add_heading("Análise da IA", level=1)
        _adicionar_blocos_word(doc, mdu.parse_blocks(resultado["narrativa_ia"]))

    buf = io.BytesIO()
    # rodapé da marca
    _rodape_word(doc)
    doc.save(buf)
    buf.seek(0)
    return buf


def gerar_pdf_planilha(resultado: dict, info: dict = None) -> io.BytesIO:
    """Gera um PDF resumindo a planilha completa."""
    info = info or resultado.get("planilha_info", {})
    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4,
                            leftMargin=15 * mm, rightMargin=15 * mm,
                            topMargin=15 * mm, bottomMargin=15 * mm)
    estilos = getSampleStyleSheet()
    titulo = ParagraphStyle("Titulo", parent=estilos["Title"], fontSize=18, textColor=HexColor("#1F3B73"))
    h2 = ParagraphStyle("H2", parent=estilos["Heading2"], fontSize=13, textColor=HexColor("#1F3B73"), spaceBefore=10)
    corpo = ParagraphStyle("Corpo", parent=estilos["BodyText"], fontSize=9, leading=12)
    nota = ParagraphStyle("Nota", parent=estilos["BodyText"], fontSize=8, textColor=HexColor("#555555"), italic=True)

    elementos = []
    _marca_pdf(doc, elementos)
    elementos.append(Paragraph("Análise da Planilha", titulo))
    elementos.append(Paragraph(f"{resultado.get('nome_arquivo','planilha')} · {info.get('n_abas',0)} abas", nota))
    elementos.append(Spacer(1, 8))

    elementos.append(Paragraph("Abas da planilha", h2))
    linhas_tab = [["Aba", "Linhas", "Colunas"]] + [
        [a.get("nome",""), str(a.get("linhas_lidas",0)), str(a.get("colunas_max",0))] for a in info.get("abas", [])
    ]
    tbl = Table(linhas_tab, colWidths=[80*mm, 35*mm, 35*mm])
    tbl.setStyle(TableStyle([
        ("FONTNAME", (0,0), (-1,0), "Helvetica-Bold"),
        ("FONTSIZE", (0,0), (-1,-1), 8),
        ("GRID", (0,0), (-1,-1), 0.4, HexColor("#CCCCCC")),
        ("BACKGROUND", (0,0), (-1,0), HexColor("#1F3B73")),
        ("TEXTCOLOR", (0,0), (-1,0), HexColor("#FFFFFF")),
    ]))
    elementos.append(tbl)
    elementos.append(Spacer(1, 8))

    for a in info.get("abas", []):
        elementos.append(Paragraph(f"Aba: {a.get('nome','')}", h2))
        for c in a.get("amostra", [])[:6]:
            elementos.append(Paragraph(" | ".join(str(x) for x in c)[:120], corpo))

    if resultado.get("narrativa_ia"):
        elementos.append(Paragraph("Análise da IA", h2))
        _adicionar_blocos_pdf(elementos, mdu.parse_blocks(resultado["narrativa_ia"]), corpo, h2)

    _rodape_pdf(elementos, "TRUST CORPORATION")
    doc.build(elementos)
    buf.seek(0)
    return buf


def gerar_xlsx_processado(resultado: dict, info: dict = None) -> io.BytesIO:
    """Gera um XLSX com o resumo estruturado da planilha (abas, conteúdo)."""
    import openpyxl as _ox

    info = info or resultado.get("planilha_info", {})
    wb = _ox.Workbook()
    ws = wb.active
    ws.title = "Resumo"
    ws.append(["Aba", "Linhas", "Colunas"])
    for a in info.get("abas", []):
        ws.append([a.get("nome",""), a.get("linhas_lidas",0), a.get("colunas_max",0)])

    # aba por aba com amostra
    for a in info.get("abas", []):
        nome_aba = str(a.get("nome",""))[:28] or "Aba"
        ws_aba = wb.create_sheet(title=nome_aba)
        for c in a.get("cabecalhos", []):
            ws_aba.append([str(x) for x in c])
        for c in a.get("amostra", []):
            ws_aba.append([str(x) for x in c])

    if resultado.get("narrativa_ia"):
        ws_ia = wb.create_sheet(title="Analise IA")
        ws_ia.append(["Seção", "Conteúdo"])
        for b in mdu.parse_blocks(resultado["narrativa_ia"]):
            tipo = b["tipo"]
            if tipo in ("h1", "h2", "h3", "h4"):
                ws_ia.append([b["texto"], ""])
            elif tipo == "p":
                ws_ia.append(["", b["texto"]])
            elif tipo in ("ul", "ol"):
                for item in b["itens"]:
                    ws_ia.append(["• item", item])
            elif tipo == "blockquote":
                ws_ia.append(["destaque", b["texto"]])
            elif tipo == "table":
                for linha in b["linhas"]:
                    ws_ia.append(["", " | ".join(str(c) for c in linha)])

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf
